from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Helper function to add a slide with title and content
def add_slide(title, content, bullet_points=None, notes=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    if bullet_points:
        tf = content_placeholder.text_frame
        tf.clear()
        for i, point in enumerate(bullet_points):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = point
            p.level = 0
    else:
        content_placeholder.text = content
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "VISHUDHH Portal: AI-Powered Cleanliness Monitoring"
slide.placeholders[1].text = "Empowering Communities & Authorities for a Cleaner Nation\n\nPresented by: [Your Name(s)]\n[Your Institution/Organization]\n[Date]"

# Slide 2: Problem Statement
add_slide(
    "Problem Statement",
    "Maintaining cleanliness in public spaces is a persistent challenge in India. Manual monitoring is inefficient, lacks transparency, and does not provide real-time feedback. There is a need for a scalable, data-driven solution that empowers both authorities and communities to monitor, manage, and improve cleanliness using modern technology."
)

# Slide 3: Team Details
add_slide(
    "Team Details",
    "",
    bullet_points=[
        "Team Name: [Your Team Name]",
        "Members:",
        "- [Member 1 Name] – [Role]",
        "- [Member 2 Name] – [Role]",
        "- [Member 3 Name] – [Role]",
        "- [Member 4 Name] – [Role]",
        "Institution/Organization: [Your College/Company]"
    ]
)

# Slide 4: Idea/Approach Description
add_slide(
    "Idea/Approach Description",
    "VISHUDHH Portal is an AI-powered web platform designed to automate cleanliness monitoring in public spaces. Using real-time camera feeds and a YOLO-based AI model, the system detects and classifies waste, calculates a cleanliness score, and presents the data on role-based dashboards for administrators at national, state, and district levels. The platform also engages the community through pledges, event discovery, and impact tracking, fostering collective action for a cleaner environment.\n\nKey Features:",
    bullet_points=[
        "Real-time waste detection and cleanliness scoring using AI",
        "Hierarchical dashboards for different administrative roles",
        "Automated email reporting and notifications",
        "Community engagement modules (pledges, events, impact)"
    ]
)

# Slide 5: Process Flowchart / Visual Representation
add_slide(
    "Process Flowchart / Visual Representation",
    "See diagram on next slide or below:",
    bullet_points=[
        "User/Camera",
        "    |",
        "    v",
        "[React Frontend] --(Image)--> [Flask AI/ML Backend]",
        "    |                                 |",
        "    |<--(Score, Detected Items)-------|",
        "    |",
        "    v",
        "[Node.js Backend] <--> [MongoDB]",
        "    |",
        "    v",
        "[Dashboards, Email Reports, Community Features]"
    ]
)

# Slide 6: Technology Stack
add_slide(
    "Technology Stack",
    "",
    bullet_points=[
        "Frontend:",
        "- React.js",
        "- Tailwind CSS",
        "- Chart.js, Framer Motion",
        "- React Webcam, Google reCAPTCHA",
        "Backend:",
        "- Node.js",
        "- Express.js",
        "- Mongoose (MongoDB ODM)",
        "- Nodemailer (Email)",
        "AI/ML Service:",
        "- Python",
        "- Flask",
        "- YOLO (You Only Look Once) for object detection",
        "- OpenCV",
        "- MongoDB (GridFS for image storage)",
        "Database:",
        "- MongoDB (Cloud/Local)",
        "Other:",
        "- RESTful APIs",
        "- CORS, Axios (API calls)"
    ]
)

# Save the presentation
prs.save("VISHUDHH_Portal_Presentation.pptx")

print("Presentation created: VISHUDHH_Portal_Presentation.pptx") 